#!/usr/bin/env python3
"""
Create enhanced PowerPoint with visual diagrams
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image, ImageDraw, ImageFont
import io
import os

def create_system_diagram():
    """Create system architecture diagram"""
    img = Image.new('RGB', (1200, 600), color='white')
    draw = ImageDraw.Draw(img)
    
    # Colors
    colors = {
        'blue': (0, 102, 204),
        'green': (51, 153, 102),
        'orange': (255, 102, 0),
        'purple': (153, 76, 153),
        'gray': (200, 200, 200)
    }
    
    # Draw boxes
    # Android App
    draw.rectangle([100, 50, 450, 150], fill=colors['green'], outline='black', width=2)
    draw.text((160, 90), "Android App", fill='white')
    
    # Arrow
    draw.line([(450, 100), (550, 100)], fill='black', width=3)
    draw.polygon([(550, 100), (530, 90), (530, 110)], fill='black')
    
    # API
    draw.rectangle([550, 50, 900, 150], fill=colors['orange'], outline='black', width=2)
    draw.text((660, 90), "REST API (PHP)", fill='white')
    
    # Arrow 2
    draw.line([(725, 150), (725, 250)], fill='black', width=3)
    draw.polygon([(725, 250), (715, 230), (735, 230)], fill='black')
    
    # Database
    draw.rectangle([550, 250, 900, 350], fill=colors['purple'], outline='black', width=2)
    draw.text((620, 290), "MySQL Database", fill='white')
    
    # Services
    draw.rectangle([100, 250, 450, 350], fill=colors['blue'], outline='black', width=2)
    draw.text((130, 290), "Firebase / ML Kit", fill='white')
    
    # Arrow 3
    draw.line([(450, 300), (550, 300)], fill='black', width=3)
    draw.polygon([(550, 300), (530, 290), (530, 310)], fill='black')
    
    # Labels
    draw.text((10, 400), "Frontend Layer", fill='black')
    draw.text((550, 400), "Backend Layer", fill='black')
    draw.text((550, 450), "Database Layer", fill='black')
    draw.text((100, 450), "Services Layer", fill='black')
    
    path = "/tmp/system_diagram.png"
    img.save(path)
    return path

def create_gps_face_diagram():
    """Create GPS and Face verification diagram"""
    img = Image.new('RGB', (1200, 600), color='white')
    draw = ImageDraw.Draw(img)
    
    colors = {
        'green': (51, 153, 102),
        'red': (204, 0, 0),
        'blue': (0, 102, 204)
    }
    
    # GPS Section
    draw.rectangle([50, 50, 550, 550], fill=colors['blue'], outline='black', width=2)
    draw.text((150, 80), "GPS Verification", fill='white')
    draw.ellipse([150, 150, 450, 450], outline='white', width=2)
    draw.text((180, 280), "1000m Radius", fill='white')
    draw.text((180, 310), "Threshold", fill='white')
    
    # Face Section
    draw.rectangle([650, 50, 1150, 550], fill=colors['green'], outline='black', width=2)
    draw.text((750, 80), "Face Detection", fill='white')
    draw.rectangle([750, 150, 1050, 450], fill='white', outline='white', width=2)
    draw.text((820, 280), "95%+ Accuracy", fill=colors['green'])
    draw.text((850, 310), "ML Kit", fill=colors['green'])
    
    path = "/tmp/verification_diagram.png"
    img.save(path)
    return path

def create_timeline_diagram():
    """Create implementation timeline"""
    img = Image.new('RGB', (1400, 400), color='white')
    draw = ImageDraw.Draw(img)
    
    colors_timeline = [
        (255, 102, 0),   # Phase 1
        (0, 102, 204),   # Phase 2
        (51, 153, 102),  # Phase 3
        (153, 76, 153)   # Phase 4
    ]
    
    phases = ["Requirements", "Design", "Development", "Testing"]
    x_positions = [150, 500, 850, 1200]
    
    for i, (phase, x, color) in enumerate(zip(phases, x_positions, colors_timeline)):
        # Circle
        draw.ellipse([x-50, 100, x+50, 200], fill=color, outline='black', width=2)
        # Text
        draw.text((x-40, 220), phase, fill='black')
        
        # Arrow
        if i < len(phases) - 1:
            draw.line([(x+50, 150), (x_positions[i+1]-50, 150)], fill='black', width=2)

    path = "/tmp/timeline_diagram.png"
    img.save(path)
    return path

def create_features_diagram():
    """Create key features visualization"""
    img = Image.new('RGB', (1400, 800), color='white')
    draw = ImageDraw.Draw(img)
    
    features = [
        ("Authentication", (100, 100), (0, 102, 204)),
        ("GPS Verification", (500, 100), (51, 153, 102)),
        ("Face Detection", (900, 100), (255, 102, 0)),
        ("Notifications", (100, 400), (153, 76, 153)),
        ("Reports", (500, 400), (204, 0, 0)),
        ("Mobile App", (900, 400), (51, 102, 153))
    ]
    
    for feature, (x, y), color in features:
        # Box
        draw.rectangle([x, y, x+300, y+150], fill=color, outline='black', width=2)
        # Text
        draw.text((x+50, y+60), feature, fill='white')
    
    path = "/tmp/features_diagram.png"
    img.save(path)
    return path

# Create diagrams
print("Creating diagrams...")
system_diagram = create_system_diagram()
verification_diagram = create_gps_face_diagram()
timeline_diagram = create_timeline_diagram()
features_diagram = create_features_diagram()
print("✅ Diagrams created")

# Load existing presentation
prs = Presentation("/Users/anshu/sams-backend/SAMS_FINAL_YEAR_PROJECT.pptx")

# Insert diagram slide after slide 4 (Architecture Overview)
# Add System Diagram Slide
blank_slide = prs.slides.add_slide(prs.slide_layouts[6])
background = blank_slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(245, 245, 245)

# Title
title_box = blank_slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
title_frame = title_box.text_frame
p = title_frame.paragraphs[0]
p.text = "System Architecture Diagram"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = RGBColor(0, 102, 204)

# Add image
img_stream = blank_slide.shapes.add_picture(system_diagram, Inches(1), Inches(1.2), width=Inches(8))

# Add Verification Slide
blank_slide2 = prs.slides.add_slide(prs.slide_layouts[6])
background2 = blank_slide2.background
fill2 = background2.fill
fill2.solid()
fill2.fore_color.rgb = RGBColor(245, 245, 245)

title_box2 = blank_slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
title_frame2 = title_box2.text_frame
p2 = title_frame2.paragraphs[0]
p2.text = "Dual Verification System"
p2.font.size = Pt(36)
p2.font.bold = True
p2.font.color.rgb = RGBColor(0, 102, 204)

img_stream2 = blank_slide2.shapes.add_picture(verification_diagram, Inches(0.5), Inches(1.2), width=Inches(9))

# Add Features Slide
blank_slide3 = prs.slides.add_slide(prs.slide_layouts[6])
background3 = blank_slide3.background
fill3 = background3.fill
fill3.solid()
fill3.fore_color.rgb = RGBColor(245, 245, 245)

title_box3 = blank_slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
title_frame3 = title_box3.text_frame
p3 = title_frame3.paragraphs[0]
p3.text = "Key Features Overview"
p3.font.size = Pt(36)
p3.font.bold = True
p3.font.color.rgb = RGBColor(0, 102, 204)

img_stream3 = blank_slide3.shapes.add_picture(features_diagram, Inches(0.2), Inches(1.1), width=Inches(9.6))

# Add Timeline Slide
blank_slide4 = prs.slides.add_slide(prs.slide_layouts[6])
background4 = blank_slide4.background
fill4 = background4.fill
fill4.solid()
fill4.fore_color.rgb = RGBColor(245, 245, 245)

title_box4 = blank_slide4.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
title_frame4 = title_box4.text_frame
p4 = title_frame4.paragraphs[0]
p4.text = "Development Timeline"
p4.font.size = Pt(36)
p4.font.bold = True
p4.font.color.rgb = RGBColor(0, 102, 204)

img_stream4 = blank_slide4.shapes.add_picture(timeline_diagram, Inches(1), Inches(1.2), width=Inches(8))

# Save enhanced presentation
output_path = "/Users/anshu/sams-backend/SAMS_FINAL_YEAR_PROJECT_Enhanced.pptx"
prs.save(output_path)
print(f"✅ Enhanced presentation saved: {output_path}")
print(f"📊 Total slides: {len(prs.slides)}")

# Clean up temp files
for f in [system_diagram, verification_diagram, timeline_diagram, features_diagram]:
    try:
        os.remove(f)
    except:
        pass
