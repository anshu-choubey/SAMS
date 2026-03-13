#!/usr/bin/env python3
"""
Create PREMIUM PowerPoint presentation with advanced graphics and visualizations
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image, ImageDraw, ImageFont
import io
import os
from math import cos, sin, pi

def create_advanced_architecture():
    """Create advanced 3D-style architecture diagram"""
    img = Image.new('RGB', (1400, 700), color=(248, 248, 248))
    draw = ImageDraw.Draw(img)
    
    # Color scheme
    colors = {
        'client': (51, 153, 102),      # Green
        'api': (255, 102, 0),          # Orange
        'db': (153, 76, 153),          # Purple
        'service': (0, 102, 204),      # Blue
        'border': (100, 100, 100),
        'shadow': (200, 200, 200)
    }
    
    # Client Layer (Android)
    x1, y1 = 100, 100
    draw.rectangle([x1, y1, x1+300, y1+150], fill=colors['client'], outline=colors['border'], width=3)
    draw.text((x1+60, y1+65), "Android Client", fill='white')
    draw.text((x1+40, y1+100), "(Kotlin Compose)", fill='white')
    
    # Arrow 1
    draw.line([(x1+300, y1+75), (x1+350, y1+75)], fill=colors['border'], width=4)
    
    # API Layer
    x2, y2 = 450, 100
    draw.rectangle([x2, y2, x2+300, y2+150], fill=colors['api'], outline=colors['border'], width=3)
    draw.text((x2+80, y2+65), "REST API", fill='white')
    draw.text((x2+65, y2+100), "(PHP 7.4+)", fill='white')
    
    # Arrow 2
    draw.line([(x2+300, y2+75), (x2+350, y2+75)], fill=colors['border'], width=4)
    
    # Database Layer
    x3, y3 = 800, 100
    draw.rectangle([x3, y3, x3+300, y3+150], fill=colors['db'], outline=colors['border'], width=3)
    draw.text((x3+60, y3+65), "MySQL Database", fill='white')
    draw.text((x3+80, y3+100), "(13 Tables)", fill='white')
    
    # Services boxes below
    services = [
        ("Firebase\nCloud Messaging", 100, 350),
        ("Google\nML Kit", 350, 350),
        ("GPS\nVerification", 600, 350),
        ("Session\nManagement", 850, 350),
        ("Encryption\nService", 1100, 350)
    ]
    
    for service, sx, sy in services:
        draw.rectangle([sx, sy, sx+250, sy+130], fill=colors['service'], outline=colors['border'], width=2)
        draw.text((sx+30, sy+40), service, fill='white')
    
    # Connection lines from API to services
    draw.line([(x2+150, y2+150), (225, 350)], fill='gray', width=2)
    draw.line([(x2+150, y2+150), (475, 350)], fill='gray', width=2)
    draw.line([(x2+150, y2+150), (725, 350)], fill='gray', width=2)
    draw.line([(x2+150, y2+150), (975, 350)], fill='gray', width=2)
    draw.line([(x2+150, y2+150), (1225, 350)], fill='gray', width=2)
    
    # Title
    draw.text((400, 550), "Complete System Architecture with Services", fill=(0, 0, 0))
    
    path = "/tmp/advanced_architecture.png"
    img.save(path)
    return path

def create_data_flow_diagram():
    """Create attendance data flow diagram"""
    img = Image.new('RGB', (1300, 700), color='white')
    draw = ImageDraw.Draw(img)
    
    colors = {
        'step1': (0, 102, 204),
        'step2': (51, 153, 102),
        'step3': (255, 102, 0),
        'step4': (153, 76, 153),
        'decision': (204, 0, 0),
        'text': (50, 50, 50)
    }
    
    # Step 1: Student Opens App
    x, y = 50, 80
    draw.rectangle([x, y, x+200, y+80], fill=colors['step1'], outline='black', width=2)
    draw.text((x+20, y+30), "Student Opens\nApp", fill='white')
    
    # Arrow
    draw.line([(x+200, y+40), (x+250, y+40)], fill='black', width=2)
    draw.polygon([(x+250, y+40), (x+240, y+30), (x+240, y+50)], fill='black')
    
    # Step 2: GPS Capture
    x = 300
    draw.rectangle([x, y, x+200, y+80], fill=colors['step2'], outline='black', width=2)
    draw.text((x+40, y+30), "Capture GPS\nLocation", fill='white')
    
    # Arrow
    draw.line([(x+200, y+40), (x+250, y+40)], fill='black', width=2)
    draw.polygon([(x+250, y+40), (x+240, y+30), (x+240, y+50)], fill='black')
    
    # Step 3: Face Detection
    x = 550
    draw.rectangle([x, y, x+200, y+80], fill=colors['step3'], outline='black', width=2)
    draw.text((x+30, y+30), "Face Detection\n(ML Kit)", fill='white')
    
    # Arrow down
    draw.line([(x+100, y+80), (x+100, y+150)], fill='black', width=2)
    draw.polygon([(x+100, y+150), (x+90, y+140), (x+110, y+140)], fill='black')
    
    # Decision Diamond
    diamond_x, diamond_y = 50, 200
    points = [
        (diamond_x+100, diamond_y),           # top
        (diamond_x+180, diamond_y+80),        # right
        (diamond_x+100, diamond_y+160),       # bottom
        (diamond_x+20, diamond_y+80)          # left
    ]
    draw.polygon(points, fill=colors['decision'], outline='black', width=2)
    draw.text((diamond_x+40, diamond_y+50), "Both Pass?", fill='white')
    
    # Success Path - Yes
    draw.line([(diamond_x+180, diamond_y+80), (diamond_x+280, diamond_y+80)], fill='green', width=3)
    draw.text((diamond_x+200, diamond_y+50), "YES", fill='green')
    
    # Success Box
    draw.rectangle([diamond_x+280, diamond_y+40, diamond_x+480, diamond_y+120], 
                   fill=(51, 153, 102), outline='black', width=2)
    draw.text((diamond_x+305, diamond_y+65), "Mark PRESENT", fill='white')
    
    # Failure Path - No
    draw.line([(diamond_x+100, diamond_y+160), (diamond_x+100, diamond_y+250)], fill='red', width=3)
    draw.text((diamond_x+110, diamond_y+200), "NO", fill='red')
    
    # Failure Box
    draw.rectangle([diamond_x, diamond_y+250, diamond_x+200, diamond_y+330], 
                   fill=(204, 0, 0), outline='black', width=2)
    draw.text((diamond_x+20, diamond_y+285), "Mark ABSENT", fill='white')
    
    # Database storage
    x, y = 750, 200
    draw.rectangle([x, y, x+250, y+100], fill=(153, 76, 153), outline='black', width=2)
    draw.text((x+40, y+35), "Store in Database\n+ Notify Server", fill='white')
    
    # Connect success to database
    draw.line([(diamond_x+480, diamond_y+80), (x, y+50)], fill='gray', width=2)
    
    # Title
    draw.text((300, 550), "Attendance Marking Data Flow", fill=colors['text'])
    
    path = "/tmp/data_flow_diagram.png"
    img.save(path)
    return path

def create_statistics_visual():
    """Create statistics and metrics visualization"""
    img = Image.new('RGB', (1400, 700), color='white')
    draw = ImageDraw.Draw(img)
    
    # Bar chart data
    stats = [
        ("API\nEndpoints", 50, (0, 102, 204)),
        ("Database\nTables", 13, (51, 153, 102)),
        ("Code\nExamples", 50, (255, 102, 0)),
        ("Test\nCases", 150, (153, 76, 153)),
        ("Documentation\nPages", 100, (204, 0, 0))
    ]
    
    max_value = max(s[1] for s in stats)
    chart_height = 400
    bar_width = 200
    start_x = 100
    start_y = 500
    
    for i, (label, value, color) in enumerate(stats):
        x = start_x + i * 250
        
        # Scale value to height
        bar_height = (value / max_value) * chart_height
        
        # Draw bar
        draw.rectangle([x, start_y - bar_height, x + bar_width, start_y], 
                       fill=color, outline='black', width=2)
        
        # Draw value on top
        draw.text((x + 60, start_y - bar_height - 30), str(value), fill='black')
        
        # Draw label
        draw.text((x + 30, start_y + 20), label, fill='black')
    
    # Title
    draw.text((350, 50), "Project Statistics & Metrics", fill='black')
    
    # Grid lines
    for h in range(0, 500, 100):
        y = start_y - h
        draw.line([(50, y), (1350, y)], fill='lightgray', width=1)
    
    path = "/tmp/statistics_visual.png"
    img.save(path)
    return path

def create_security_layers():
    """Create security layers visualization"""
    img = Image.new('RGB', (1400, 700), color='white')
    draw = ImageDraw.Draw(img)
    
    layers = [
        ("HTTPS/TLS", 150, 515, (0, 102, 204)),
        ("Input Validation", 150, 415, (51, 153, 102)),
        ("Authentication", 150, 315, (255, 102, 0)),
        ("Encryption (AES-256)", 150, 215, (153, 76, 153)),
        ("Database Security", 150, 115, (204, 0, 0))
    ]
    
    # Draw layers
    for label, x, y, color in layers:
        draw.rectangle([x, y, x+400, y+80], fill=color, outline='black', width=2)
        draw.text((x+50, y+30), label, fill='white')
    
    # Add descriptions
    descriptions = [
        "Secure Communication",
        "Prevent XSS, SQL Injection",
        "Session-based with bcrypt",
        "Sensitive Data Protection",
        "Prepared Statements"
    ]
    
    for i, (desc, x, y, _) in enumerate(layers):
        draw.text((x+430, y+30), descriptions[i], fill='black')
    
    # Title
    draw.text((300, 50), "Multi-Layer Security Architecture", fill='black')
    
    # Shield icon
    draw.polygon([(1200, 100), (1300, 100), (1300, 300), (1250, 350), (1200, 300)], 
                 fill=(0, 102, 204), outline='black', width=2)
    
    path = "/tmp/security_layers.png"
    img.save(path)
    return path

def create_comparison_table():
    """Create manual vs automated attendance comparison"""
    img = Image.new('RGB', (1400, 600), color='white')
    draw = ImageDraw.Draw(img)
    
    # Header
    draw.rectangle([0, 0, 1400, 80], fill=(0, 102, 204), outline='black', width=2)
    draw.text((300, 25), "Manual Attendance vs Automated SAMS", fill='white')
    
    # Column headers
    draw.rectangle([50, 80, 650, 150], fill=(51, 153, 102), outline='black', width=1)
    draw.text((200, 100), "Manual System", fill='white')
    
    draw.rectangle([650, 80, 1350, 150], fill=(255, 102, 0), outline='black', width=1)
    draw.text((950, 100), "SAMS (Automated)", fill='white')
    
    # Comparison rows
    comparisons = [
        ("Time per class", "30-45 minutes", "2-3 minutes"),
        ("Accuracy", "70-80% (human error)", "99%+ (GPS + Face)"),
        ("Proxying", "High risk", "Prevented"),
        ("Real-time tracking", "Manual reporting", "Instant updates"),
        ("Cost per term", "$500+ (staff hours)", "$50 (server)"),
        ("Data security", "Physical sheets", "Encrypted database"),
        ("Parent notification", "Weekly reports", "Real-time alerts"),
    ]
    
    y = 150
    for metric, manual, automated in comparisons:
        # Metric column
        draw.rectangle([50, y, 650, y+60], fill=(240, 240, 240), outline='gray', width=1)
        draw.text((70, y+20), metric, fill='black')
        
        # Manual column
        draw.rectangle([650, y, 1350, y+60], fill=(255, 240, 220), outline='gray', width=1)
        draw.text((670, y+10), f"Manual: {manual}", fill='black')
        draw.text((670, y+30), f"SAMS: {automated}", fill=(0, 153, 0))
        
        y += 60
    
    path = "/tmp/comparison_table.png"
    img.save(path)
    return path

print("🎨 Creating premium diagrams...")

# Create all diagrams
advanced_arch = create_advanced_architecture()
data_flow = create_data_flow_diagram()
stats = create_statistics_visual()
security = create_security_layers()
comparison = create_comparison_table()

print("✅ All diagrams created successfully")

# Load existing enhanced presentation
prs = Presentation("/Users/anshu/sams-backend/SAMS_FINAL_YEAR_PROJECT_Enhanced.pptx")

def add_image_slide(title, image_path):
    """Add slide with image"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 245, 245)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9.4), Inches(0.6))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 102, 204)
    
    # Add image
    slide.shapes.add_picture(image_path, Inches(0.3), Inches(0.9), width=Inches(9.4))
    
    return slide

# Add new slides with diagrams
add_image_slide("Advanced System Architecture", advanced_arch)
add_image_slide("Attendance Data Flow & Processing", data_flow)
add_image_slide("Security Layers & Protection", security)
add_image_slide("Project Metrics & Statistics", stats)
add_image_slide("Manual vs Automated Comparison", comparison)

# Add additional content slides

# Slide: Key Performance Indicators
slide = prs.slides.add_slide(prs.slide_layouts[6])
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(245, 245, 245)

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
title_frame = title_box.text_frame
p = title_frame.paragraphs[0]
p.text = "Key Performance Indicators (KPIs)"
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = RGBColor(0, 102, 204)

# KPI boxes
kpis = [
    ("Response Time", "<500ms", RGBColor(0, 102, 204)),
    ("Face Accuracy", "99%+", RGBColor(51, 153, 102)),
    ("GPS Accuracy", "95%", RGBColor(255, 102, 0)),
    ("Uptime", "99.5%", RGBColor(153, 76, 153)),
    ("Concurrent Users", "10,000+", RGBColor(204, 0, 0)),
    ("Data Sync", "<2 seconds", RGBColor(0, 153, 102))
]

x_positions = [0.5, 3.5, 6.5]
y_positions = [1.5, 4]

for idx, (metric, value, color) in enumerate(kpis):
    x = x_positions[idx % 3]
    y = y_positions[idx // 3]
    
    # Box
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y),
        Inches(2.8), Inches(2)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color
    
    # Metric text
    text_frame = shape.text_frame
    text_frame.word_wrap = True
    
    p = text_frame.paragraphs[0]
    p.text = metric
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    p = text_frame.add_paragraph()
    p.text = value
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

# Slide: Technology Comparison
slide = prs.slides.add_slide(prs.slide_layouts[6])
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(245, 245, 245)

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
title_frame = title_box.text_frame
p = title_frame.paragraphs[0]
p.text = "Technology Choices: Why These?"
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = RGBColor(0, 102, 204)

tech_choices = [
    ("PHP 7.4+", "Fast, reliable, excellent MySQL support", 0.5, 1.5),
    ("Kotlin", "Type-safe, concise, 100% interoperable with Java", 0.5, 3),
    ("Jetpack Compose", "Modern declarative UI, hot reload, reactive", 5.5, 1.5),
    ("Firebase", "Real-time messaging, authentication, analytics", 5.5, 3),
    ("MySQL", "Scalable, ACID compliance, excellent indexing", 0.5, 4.5),
    ("ML Kit", "On-device processing, 95%+ accuracy, privacy", 5.5, 4.5)
]

for tech, reason, x, y in tech_choices:
    # Box
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(4), Inches(1.2))
    text_frame = box.text_frame
    text_frame.word_wrap = True
    
    p = text_frame.paragraphs[0]
    p.text = f"✓ {tech}"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 102, 204)
    
    p = text_frame.add_paragraph()
    p.text = reason
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(50, 50, 50)

# Slide: Implementation Highlights
slide = prs.slides.add_slide(prs.slide_layouts[6])
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(245, 245, 245)

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
title_frame = title_box.text_frame
p = title_frame.paragraphs[0]
p.text = "Implementation Highlights"
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = RGBColor(0, 102, 204)

highlights = [
    ("Code Quality", ["50+ working code examples", "80%+ test coverage", "Production-ready patterns"], RGBColor(0, 102, 204)),
    ("Security", ["AES-256 encryption", "Bcrypt hashing (cost 12)", "SQL injection prevention"], RGBColor(51, 153, 102)),
    ("Scalability", ["Horizontal scaling ready", "Database indexing optimized", "Stateless API design"], RGBColor(255, 102, 0)),
    ("Documentation", ["Complete API docs", "Architecture guides", "Deployment instructions"], RGBColor(153, 76, 153))
]

x_pos = [0.5, 5.2]
y_pos = [1.5, 4]

for idx, (category, items, color) in enumerate(highlights):
    x = x_pos[idx % 2]
    y = y_pos[idx // 2]
    
    # Title box
    title_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y),
        Inches(4.5), Inches(0.4)
    )
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = color
    title_shape.line.color.rgb = color
    
    title_text = title_shape.text_frame
    p = title_text.paragraphs[0]
    p.text = category
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # Items box
    items_box = slide.shapes.add_textbox(Inches(x), Inches(y+0.5), Inches(4.5), Inches(2.2))
    items_frame = items_box.text_frame
    items_frame.word_wrap = True
    
    for item in items:
        p = items_frame.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_COLOR
        p.space_before = Pt(4)
        p.space_after = Pt(4)

# Save premium presentation
output_path = "/Users/anshu/sams-backend/SAMS_FINAL_YEAR_PROJECT_PREMIUM.pptx"
prs.save(output_path)
print(f"✅ Premium presentation created: {output_path}")
print(f"📊 Total slides: {len(prs.slides)}")

# Clean up
for f in [advanced_arch, data_flow, stats, security, comparison]:
    try:
        os.remove(f)
    except:
        pass

# Define color for enhanced version
TEXT_COLOR = RGBColor(40, 40, 40)

print("✨ Presentation is now PREMIUM quality!")
